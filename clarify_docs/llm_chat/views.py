import os
import tempfile
import docx2txt
from PyPDF2 import PdfReader
from django.contrib.auth.decorators import login_required
from langchain.callbacks import get_openai_callback

from langchain.text_splitter import CharacterTextSplitter, RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from django.shortcuts import render, redirect, get_object_or_404
from django.http import JsonResponse
import openai
from dotenv import load_dotenv
from langchain.vectorstores.faiss import FAISS
from pptx import Presentation

from .forms import PDFUploadForm, PDFUpdateForm, PDFDocumentForm2
from .models import ChatMessage, PDFDocument, UserData
from users.models import Avatar
# CHUNK_SIZE = 1024
# CHUNK_OVERLAP = 200
# MAX_LEN = 512
# TEMPERATURE = 0.5

MAX_SIZE_FILE = 20 * 1024 * 1024  # 20 Mb
load_dotenv()


def main(request):
    avatar = Avatar.objects.filter(user_id=request.user.id).first()

    return render(request, 'llm_chat/index.html', context={'avatar': avatar,})  # home.html


def get_pdf_text(file):
    text = None
    if file:
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(file.read())
            temp_file.flush()
            if file.name.endswith('.pdf'):
                pdf_reader = PdfReader(temp_file.name)
                text = ''.join(page.extract_text() for page in pdf_reader.pages)
            elif file.name.endswith('.txt'):
                with open(temp_file.name, 'r') as f:
                    text = f.read()
            elif file.name.endswith('.docx'):
                text = docx2txt.process(temp_file.name)
            elif file.name.endswith('.pptx'):
                prs = Presentation(temp_file.name)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text = '\n'.join(text_runs)
    return text


def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore(text_chunks):
    api_key = os.getenv("OPENAI_API_KEY")
    embeddings = OpenAIEmbeddings(model="text-embedding-ada-002", openai_api_key=api_key)
    # embeddings = OllamaEmbeddings(model_name=EMBEDDING_MODEL)
    # embeddings = HuggingFaceInstructEmbeddings(
    #     model_name="hkunlp/instructor-large", model_kwargs={"device": 'cpu'}
    # )
    # embeddings = HuggingFaceInstructEmbeddings(model_name=EMBEDDING_MODEL)
    vectorstore = FAISS.from_texts(text_chunks, embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    llm = ChatOpenAI(model="gpt-3.5-turbo", temperature=0.7)
    # llm = AutoModelForCausalLM.from_pretrained("TheBloke/Llama-2-7b-Chat-GGUF",
    #                                            model_file="llama-2-7b-chat.q4_K_M.gguf", model_type="llama",
    #                                            gpu_layers=0, from_tf=True)
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain


@login_required(login_url="/login/")
def upload_pdf(request):
    
    user = request.user
    avatar = Avatar.objects.filter(user_id=user.id).first()
    try:
        user_data = UserData.objects.get(user=user)
    except UserData.DoesNotExist:
        user_data = UserData.objects.create(user=user, total_files_uploaded=0, total_questions_asked=0)

    if request.method == 'POST':  # and request.FILES.get('pdf_document'):
        form = PDFUploadForm(request.POST, request.FILES)
        if form.is_valid():
            pdf_document = request.FILES['pdf_document']

            _, file_extension = os.path.splitext(pdf_document.name)
            if file_extension.lower() not in [".pdf", ".txt", ".docx", ".pptx"]:
                return JsonResponse({'error': 'Only PDF, TXT, DOCX, PPTX files'}, status=400)

            if pdf_document.size > MAX_SIZE_FILE:  # 20 Mb
                return JsonResponse({'error': "File size exceeds 50 MB."}, status=400)

            # Check if a file with the same name already exists for this user
            if PDFDocument.objects.filter(user=user, title=pdf_document.name).exists():
                return JsonResponse({'error': 'A file with the same name already exists.'}, status=400)
            user_data.total_files_uploaded += 1
            user_data.save()
            # Save PDF -> DB
            pdf = PDFDocument(user=user, title=pdf_document.name)
            pdf.documentContent = get_pdf_text(pdf_document)
            pdf.save()
            return JsonResponse({'success': 'File uploaded successfully.'})

    else:
        form = PDFUploadForm()
    user_pdfs = PDFDocument.objects.filter(user=request.user)
    # not needed for now
    chat_message = ChatMessage.objects.all()
    return render(request, 'llm_chat/chat.html', {'form': form, 'user_pdfs': user_pdfs, 'avatar': avatar})


@login_required(login_url="/login/")
def ask_question(request):
    user = request.user
    avatar = Avatar.objects.filter(user_id=user.id).first()
    try:
        user_data = UserData.objects.get(user=user)
    except UserData.DoesNotExist:
        user_data = UserData.objects.create(user=user, total_files_uploaded=0, total_questions_asked=0)

    chat_history = ChatMessage.objects.filter(user=request.user).order_by('timestamp')[:10]
    chat_response = ''
    user_pdfs = PDFDocument.objects.filter(user=request.user)
    user_question = ""
    selected_pdf = None  # selected_pdf_id або об'єкт PDFDocument, треба перевіряти ще

    if request.method == 'POST':
        # Check if the user has exceeded the limit for questions per file
        user_question = request.POST.get('user_question')
        selected_pdf_id = request.POST.get('selected_pdf')
        selected_pdf = get_object_or_404(PDFDocument, id=selected_pdf_id)
        text_chunks = get_text_chunks(selected_pdf.documentContent)

        knowledge_base = get_vectorstore(text_chunks)
        conversation_chain = get_conversation_chain(knowledge_base)

        with get_openai_callback() as cb:
            response = conversation_chain({'question': user_question})

        chat_response = response["answer"]
        chat_message = ChatMessage(user=request.user, message=user_question, answer=chat_response,
                                   pdf_document=selected_pdf)  # об'єкт PDFDocument
        user_data.total_questions_asked += 1
        user_data.save()
        chat_message.save()

    # Отримуємо повідомлення, які відносяться до обраного PDFDocument
    chat_message = ChatMessage.objects.filter(user=request.user, pdf_document=selected_pdf).order_by('timestamp')

    context = {'chat_response': chat_response, 'chat_history': chat_history, 'user_question': user_question,
               'user_pdfs': user_pdfs, 'chat_message': chat_message, 'avatar': avatar}

    return render(request, 'llm_chat/chat.html', context)


@login_required(login_url="/login/")
def get_chat_history(request):
    pdf_id = request.GET.get('pdf_id')
    if not pdf_id:
        return JsonResponse({"error": "PDF ID not provided"}, status=400)

    # Тут історія чату для даного пдф, дивитись фронт зда'ться там проблема
    messages = ChatMessage.objects.filter(pdf_document_id=pdf_id).values('message', 'answer', 'timestamp')

    return JsonResponse(list(messages), safe=False)


# @login_required(login_url="/login/")
# def view_pdf(request, pdf_id):
#     """
#     Displays a PDF document.

#     :param request: HTTP request.
#     :param pdf_id: ID of the PDF document.
#     :return: Rendered page with the PDF document.
#     """
#     pdf = PDFDocument.objects.get(id=pdf_id)
#     return render(request, 'view_pdf.html', {'pdf': pdf})

# @login_required(login_url="/login/")
# def list_pdfs(request):
#     """
#     Lists PDF documents for the logged-in user.

#     :param request: HTTP request.
#     :return: Rendered page with the list of PDF documents.
#     """
#     pdfs = PDFDocument.objects.filter(user=request.user)
#     return render(request, 'edit_pdf.html', {'pdfs': pdfs})

# @login_required(login_url="/login/")
# def delete_pdf(request, pdf_id):
#     """
#     Deletes a PDF document.

#     :param request: HTTP request.
#     :param pdf_id: ID of the PDF document to delete.
#     :return: Redirect or error response.
#     """
#     try:
#         pdfs = PDFDocument.objects.get(id=pdf_id)
#         pdf_title = getattr(pdfs, 'document', pdfs.title)
#         pdfs.delete()
#         return redirect('/pdfs/')
#     except PDFDocument.DoesNotExist:
#         # Handle the case where the PDFDocument with the given id does not exist
#         return HttpResponseNotFound("PDF not found")

# @login_required(login_url="/login/")
# def update_pdf(request, pdf_id):
#     """
#     Updates a PDF document.

#     :param request: HTTP request.
#     :param pdf_id: ID of the PDF document to update.
#     :return: Rendered page with the updated PDF document.
#     """
#     pdf = get_object_or_404(PDFDocument, pk=pdf_id)
#     if request.method == 'POST':
#         form = PDFDocumentForm2(request.POST, instance=pdf)
#         if form.is_valid():
#             form.save()
#             messages.success(request, 'Pdf updated successfully.')
#             return redirect('/pdfs/')
#     else:
#         form = PDFUpdateForm(instance=pdf)
#     return render(request, 'update_pdf.html', {'form': form, 'pdf': pdf})

# TO DO for future
# def new_chat(request):
# clear the messages list
# request.session.pop('messages', None)
# return redirect('home')


# this is the view for handling errors
# def error_handler(request):
#     return render(request, 'llm_chat/404.html')
